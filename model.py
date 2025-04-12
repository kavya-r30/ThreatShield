from langchain_groq import ChatGroq
from langchain.prompts import ChatPromptTemplate
from langchain.schema import StrOutputParser
import json
import os
import mimetypes
import base64
import re
import olefile
import hashlib
import magic
import pdfplumber
import io
import zipfile
import struct
import yara
import math
import docx
import binascii
import logging

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)
logging.getLogger('pdfminer.pdfpage').setLevel(logging.ERROR)


os.environ["GROQ_API_KEY"] = "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"

model = ChatGroq(model="qwen-2.5-32b")

# Define YARA rules for detecting suspicious patterns
YARA_RULES = """
rule SuspiciousBase64
{
    strings:
        $b64_pattern = /[A-Za-z0-9+\/]{50,}={0,2}/
    condition:
        $b64_pattern
}

rule SuspiciousShellCommands
{
    strings:
        $cmd1 = "powershell -enc" nocase
        $cmd2 = "cmd.exe /c" nocase
        $cmd3 = "wget http" nocase
        $cmd4 = "curl http" nocase
        $cmd5 = "Invoke-Expression" nocase
        $cmd6 = "iex(" nocase
    condition:
        any of them
}

rule SuspiciousJavascript
{
    strings:
        $eval = "eval(" nocase
        $doc_write = "document.write(" nocase
        $fromcharcode = "fromCharCode" nocase
        $unescape = "unescape(" nocase
    condition:
        any of them
}
"""

# Compile YARA rules
try:
    yara_rules = yara.compile(source=YARA_RULES)
except:
    logger.warning("YARA not available, skipping YARA rule compilation")
    yara_rules = None

def get_file_hashes(file_path):
    """Calculate MD5, SHA-1 and SHA-256 hashes for a file."""
    with open(file_path, 'rb') as f:
        content = f.read()
        md5 = hashlib.md5(content).hexdigest()
        sha1 = hashlib.sha1(content).hexdigest()
        sha256 = hashlib.sha256(content).hexdigest()
    return {
        "md5": md5,
        "sha1": sha1,
        "sha256": sha256
    }

def get_file_entropy(data):
    """Calculate Shannon entropy of data to detect potential encryption or obfuscation."""
    if not data:
        return 0
    entropy = 0
    for x in range(256):
        p_x = float(data.count(bytes([x]))) / len(data)
        if p_x > 0:
            entropy += - p_x * math.log(p_x, 2)
    return entropy

def detect_embedded_files(file_path):
    """Detect potentially embedded files using magic numbers."""
    with open(file_path, 'rb') as f:
        content = f.read()
    
    embedded_types = []
    # Check for common file signatures
    signatures = {
        b'MZ': 'PE/EXE file',
        b'PK\x03\x04': 'ZIP/Office file',
        b'\x50\x4B\x05\x06': 'ZIP/Office file (empty)',
        b'\x50\x4B\x07\x08': 'ZIP/Office file (spanned)',
        b'%PDF': 'PDF file',
        b'\xFF\xD8\xFF': 'JPEG image',
        b'\x89PNG': 'PNG image',
        b'GIF87a': 'GIF image',
        b'GIF89a': 'GIF image',
        b'\x7FELF': 'ELF file',
        b'#!': 'Shell script',
        b'\x1F\x8B\x08': 'GZIP archive',
        b'BZh': 'BZIP2 archive',
        b'RIFF': 'RIFF container (AVI/WAV)',
        b'\x4D\x5A': 'Windows executable'
    }
    
    # Search for signatures after the first 4KB (might indicate embedded files)
    if len(content) > 4096:
        for offset in range(4096, len(content) - 20, 512):  # Check every 512 bytes
            chunk = content[offset:offset+20]
            for sig, file_type in signatures.items():
                if chunk.startswith(sig):
                    embedded_types.append({
                        "type": file_type,
                        "offset": offset,
                        "signature": binascii.hexlify(sig).decode('ascii')
                    })
    
    return embedded_types

def extract_pdf_content(file_path):
    """Extract detailed content and metadata from PDF files."""
    results = {
        "file_type": "pdf",
        "extension": ".pdf",
        "size_bytes": os.path.getsize(file_path),
        "metadata": {},
        "indicators": []
    }
    
    try:
        with pdfplumber.open(file_path) as pdf:
            # Extract metadata
            results["metadata"] = {
                "page_count": len(pdf.pages),
                "pdf_info": pdf.metadata
            }
            
            # Check for JavaScript
            has_js = False
            js_content = []
            
            # Extract text sample from first few pages
            text_sample = []
            for i, page in enumerate(pdf.pages):
                if i < 3:  # First 3 pages
                    try:
                        text = page.extract_text()
                        if text:
                            text_sample.append(text)
                    except Exception as e:
                        results["indicators"].append({
                            "type": "pdf_error",
                            "details": f"Error extracting text from page {i+1}: {str(e)}"
                        })
            
            # Check for suspicious patterns in text
            suspicious_patterns = [
                "javascript:", "eval(", "unescape(", "String.fromCharCode(",
                "ActiveXObject", "getAnnots", "this.submitForm", "util.exec",
                "app.alert", "ren.executeMenuItem", "runtime.exec("
            ]
            
            full_text = " ".join(text_sample)
            for pattern in suspicious_patterns:
                if pattern.lower() in full_text.lower():
                    results["indicators"].append({
                        "type": "suspicious_pdf_content",
                        "details": f"Found suspicious pattern '{pattern}' in PDF text"
                    })
            
            # Look for forms and JavaScript
            form_fields = []
            for page in pdf.pages[:5]:  # Check first 5 pages for forms
                annots = page.annots
                if annots:
                    for annot in annots:
                        if isinstance(annot, dict):
                            if annot.get('Subtype') == 'Widget':
                                form_fields.append(annot)
                            if '/JavaScript' in str(annot) or '/JS' in str(annot):
                                has_js = True
                                js_content.append(str(annot))
            
            results["metadata"]["has_form_fields"] = len(form_fields) > 0
            results["metadata"]["form_field_count"] = len(form_fields)
            results["metadata"]["has_javascript"] = has_js
            
            if has_js:
                results["indicators"].append({
                    "type": "pdf_javascript",
                    "details": "PDF contains JavaScript which may execute arbitrary code",
                    "sample": js_content[:500] if js_content else ""
                })
            
            # Check for embedded files
            try:
                embedded_files = []
                catalog = pdf.doc.catalog
                if hasattr(catalog, 'get') and catalog.get('Names'):
                    names = catalog.get('Names')
                    if names and names.get('EmbeddedFiles'):
                        embedded_files = ["PDF has embedded files"]
                
                results["metadata"]["has_embedded_files"] = len(embedded_files) > 0
                
                if embedded_files:
                    results["indicators"].append({
                        "type": "pdf_embedded_files",
                        "details": "PDF contains embedded files which may include malicious content"
                    })
            except:
                pass
                
            # Content sample
            results["content_sample"] = "\n".join(text_sample[:3])
            
    except Exception as e:
        results["extraction_error"] = str(e)
        results["indicators"].append({
            "type": "pdf_extraction_error",
            "details": f"Error analyzing PDF: {str(e)}"
        })
    
    return results

def extract_office_document_content(file_path):
    """Extract content and metadata from Office documents using specialized libraries."""
    extension = os.path.splitext(file_path)[1].lower()
    file_size = os.path.getsize(file_path)
    results = {
        "file_type": "office_document",
        "extension": extension,
        "size_bytes": file_size,
        "metadata": {},
        "indicators": []
    }
    
    try:
        # For DOCX files
        if extension == '.docx':
            import docx
            doc = docx.Document(file_path)
            
            # Extract text content
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            
            # Extract metadata
            core_properties = doc.core_properties
            results["metadata"] = {
                "author": core_properties.author,
                "created": str(core_properties.created) if core_properties.created else None,
                "last_modified_by": core_properties.last_modified_by,
                "modified": str(core_properties.modified) if core_properties.modified else None,
                "title": core_properties.title,
                "paragraph_count": len(doc.paragraphs)
            }
            
            results["content_sample"] = "\n".join(full_text[:20])  # First 20 paragraphs
            
            # Check for embedded objects and links
            try:
                # Check DOCX as ZIP file
                with zipfile.ZipFile(file_path) as zip_ref:
                    file_list = zip_ref.namelist()
                    
                    # Check for embedded objects
                    embedded_objects = [f for f in file_list if f.startswith('word/embeddings/')]
                    results["metadata"]["embedded_objects"] = embedded_objects
                    
                    # Check for external links in document.xml
                    if 'word/document.xml' in file_list:
                        doc_xml = zip_ref.read('word/document.xml').decode('utf-8', errors='ignore')
                        external_links = re.findall(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', doc_xml)
                        results["metadata"]["external_links"] = external_links
                        
                        # Check for suspicious patterns
                        suspicious_patterns = ["ActiveX", "DDEAUTO", "execute", "macro", "VBA"]
                        for pattern in suspicious_patterns:
                            if pattern.lower() in doc_xml.lower():
                                results["indicators"].append({
                                    "type": "suspicious_docx_content",
                                    "details": f"Found suspicious pattern '{pattern}' in document XML"
                                })
            except Exception as e:
                results["indicators"].append({
                    "type": "docx_analysis_error",
                    "details": f"Error analyzing DOCX structure: {str(e)}"
                })
            
        # For XLSX files
        elif extension == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, keep_vba=True)
            
            # Check for macros
            has_macros = wb.vba_archive is not None
            
            # Extract basic sheet data
            sheet_names = wb.sheetnames
            sheet_count = len(sheet_names)
            
            # Extract sample content from first sheet
            if sheet_names:
                ws = wb[sheet_names[0]]
                sample_rows = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i < 10:  # First 10 rows
                        sample_rows.append(str(row))
                    else:
                        break
            
            results["metadata"] = {
                "sheet_count": sheet_count,
                "sheet_names": sheet_names,
                "has_vba_macros": has_macros
            }
            
            if has_macros:
                results["indicators"].append({
                    "type": "excel_macros",
                    "details": "Excel file contains VBA macros which can execute code"
                })
                
            results["content_sample"] = "\n".join(sample_rows) if sheet_names else ""
            
            # Analyze for formula-based threats
            try:
                with zipfile.ZipFile(file_path) as zip_ref:
                    sheet_files = [f for f in zip_ref.namelist() if f.startswith('xl/worksheets/sheet')]
                    for sheet_file in sheet_files[:3]:  # Analyze first 3 sheets
                        sheet_content = zip_ref.read(sheet_file).decode('utf-8', errors='ignore')
                        # Check for suspicious Excel formulas
                        suspicious_formulas = ["EXEC", "CALL", "REGISTER", "=CMD", "=SHELL", 
                                              "=FORMULA", "HYPERLINK", "=DDE", "=DDEAUTO"]
                        for formula in suspicious_formulas:
                            if formula in sheet_content:
                                results["indicators"].append({
                                    "type": "suspicious_excel_formula",
                                    "details": f"Found potentially dangerous formula: {formula}"
                                })
            except Exception as e:
                results["indicators"].append({
                    "type": "xlsx_analysis_error",
                    "details": f"Error analyzing XLSX formulas: {str(e)}"
                })
                
        # For PPTX files
        elif extension == '.pptx':
            import pptx
            presentation = pptx.Presentation(file_path)
            
            # Extract text from slides
            slide_texts = []
            for i, slide in enumerate(presentation.slides):
                if i < 5:  # First 5 slides only
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            slide_text.append(shape.text)
                    slide_texts.append(" | ".join(slide_text))
            
            results["metadata"] = {
                "slide_count": len(presentation.slides)
            }
            
            # Check for embedded objects and links
            try:
                with zipfile.ZipFile(file_path) as zip_ref:
                    file_list = zip_ref.namelist()
                    # Check for embedded objects
                    embedded_objects = [f for f in file_list if '/embeddings/' in f]
                    results["metadata"]["embedded_objects"] = embedded_objects
                    
                    if embedded_objects:
                        results["indicators"].append({
                            "type": "pptx_embedded_objects",
                            "details": f"PowerPoint file contains {len(embedded_objects)} embedded objects"
                        })
                        
                    # Check for external links in slides
                    links = []
                    for f in file_list:
                        if f.startswith('ppt/slides/slide'):
                            try:
                                slide_xml = zip_ref.read(f).decode('utf-8', errors='ignore')
                                link_matches = re.findall(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', slide_xml)
                                links.extend(link_matches)
                            except:
                                pass
                    
                    results["metadata"]["external_links"] = links
            except Exception as e:
                results["indicators"].append({
                    "type": "pptx_analysis_error",
                    "details": f"Error analyzing PPTX structure: {str(e)}"
                })
                
            results["content_sample"] = "\n".join(slide_texts)
            
        # For legacy Office files (.doc, .xls, .ppt)
        elif extension in ['.doc', '.xls', '.ppt']:
            if olefile.isOleFile(file_path):
                ole = olefile.OleFileIO(file_path)
                
                # Check for macros
                has_macros = False
                vba_project = False
                for stream in ole.listdir():
                    stream_path = "/".join(stream)
                    if "macros" in stream_path.lower() or "vba" in stream_path.lower():
                        has_macros = True
                    if stream_path == "VBA/PROJECT":
                        vba_project = True
                
                # Extract metadata
                metadata = {}
                if ole.exists('\x05DocumentSummaryInformation'):
                    suminfo = ole.getproperties('\x05DocumentSummaryInformation')
                    metadata.update({str(k): str(v) for k, v in suminfo.items() if isinstance(k, (int, str))})
                if ole.exists('\x05SummaryInformation'):
                    suminfo = ole.getproperties('\x05SummaryInformation')
                    metadata.update({str(k): str(v) for k, v in suminfo.items() if isinstance(k, (int, str))})
                
                results["metadata"] = {
                    "has_vba_macros": has_macros,
                    "has_vba_project": vba_project,
                    "ole_streams": ["/".join(s) for s in ole.listdir()],
                    "extracted_properties": metadata
                }
                
                # Check for suspicious streams/markers
                suspicious_streams = []
                for stream in ole.listdir():
                    stream_path = "/".join(stream)
                    if stream_path in ["ObjectPool", "ObjInfo", "Ole10Native"]:
                        suspicious_streams.append(stream_path)
                
                if has_macros:
                    results["indicators"].append({
                        "type": "legacy_office_macros",
                        "details": "Legacy Office file contains VBA macros which can execute code"
                    })
                
                if suspicious_streams:
                    results["indicators"].append({
                        "type": "suspicious_ole_streams",
                        "details": f"Found potentially suspicious OLE streams: {', '.join(suspicious_streams)}"
                    })
                
                # We can't easily extract text content from legacy formats without additional libraries
                results["content_sample"] = "Legacy Office format - text extraction limited"
            else:
                results["metadata"] = {"error": "Not a valid OLE file"}
                results["indicators"].append({
                    "type": "invalid_office_format",
                    "details": "File has Office extension but is not a valid OLE file"
                })
        
        # For RTF files
        elif extension == '.rtf':
            with open(file_path, 'rb') as f:
                content = f.read(20000)  # Read first 20KB
                
            # Try to decode as text
            try:
                content_text = content.decode('utf-8', errors='replace')
            except:
                content_text = content.decode('latin-1', errors='replace')
                
            # Look for suspicious patterns in RTF
            suspicious_patterns = [
                "\\objocx", "\\objdata", "\\objupdate", "\\objemb", 
                "\\objautlink", "\\objclass", "\\objw", "\\objh",
                "\\datastore", "\\datafield", "\\fonttbl", "\\colortbl",
                "\\pntext", "\\listtable", "\\listoverride", "\\generator"
            ]
            
            # Object linking patterns
            obj_patterns = ["\\object", "\\objdata", "\\objupdate", "\\objemb", "\\objautlink"]
            has_obj = any(p in content_text for p in obj_patterns)
            
            # Check for hex-encoded data which might be shellcode
            has_hex_data = re.search(r'\\[0-9a-fA-F]{2}\\[0-9a-fA-F]{2}\\[0-9a-fA-F]{2}\\[0-9a-fA-F]{2}', content_text) is not None
            
            # Check for control words
            control_words = []
            for pattern in suspicious_patterns:
                if pattern in content_text:
                    control_words.append(pattern)
            
            results["metadata"] = {
                "rtf_control_words": control_words,
                "has_embedded_objects": has_obj,
                "has_hex_data": has_hex_data
            }
            
            if has_obj:
                results["indicators"].append({
                    "type": "rtf_embedded_objects",
                    "details": "RTF file contains embedded objects which may execute code"
                })
                
            if has_hex_data:
                results["indicators"].append({
                    "type": "rtf_hex_data",
                    "details": "RTF file contains hex-encoded data which may be shellcode or exploits"
                })
                
            # CVE pattern detection
            cve_patterns = {
                "CVE-2017-11882": ["\\objupdate", "\\objdata", "equation.3"],
                "CVE-2018-0802": ["\\objupdate", "\\objdata", "equation.3"],
                "CVE-2017-0199": ["urlmon.dll", "mshtml.dll", "\\objdata", "\\objw", "\\objh"],
                "CVE-2018-8174": ["VBScript.Encode", "mshtml.dll", "kernel32.dll"]
            }
            
            for cve, patterns in cve_patterns.items():
                if all(p.lower() in content_text.lower() for p in patterns):
                    results["indicators"].append({
                        "type": "rtf_exploit_pattern",
                        "details": f"RTF file contains patterns matching known exploit: {cve}"
                    })
            
            results["content_sample"] = content_text[:2000]  # First 2000 chars
            
    except ImportError as e:
        # Handle missing library dependencies
        results["metadata"] = {"error": f"Required library not installed: {str(e)}"}
        results["extraction_error"] = f"Required library not installed: {str(e)}"
    except Exception as e:
        # Handle other errors during extraction
        results["metadata"] = {"error": f"Error extracting content: {str(e)}"}
        results["extraction_error"] = str(e)
    
    return results

def extract_executable_info(file_path):
    """Extract detailed information from executable files."""
    extension = os.path.splitext(file_path)[1].lower()
    file_size = os.path.getsize(file_path)
    
    results = {
        "file_type": "executable",
        "extension": extension,
        "size_bytes": file_size,
        "metadata": {},
        "indicators": []
    }
    
    try:
        with open(file_path, 'rb') as f:
            # Read the first 4KB and last 4KB for analysis
            header = f.read(4096)
            f.seek(max(0, file_size - 4096))
            footer = f.read(4096)
        
        # Calculate entropy to detect packed/encrypted code
        entropy = get_file_entropy(header + footer)
        results["metadata"]["entropy"] = entropy
        
        if entropy > 7.0:
            results["indicators"].append({
                "type": "high_entropy",
                "details": f"File has high entropy ({entropy:.2f}/8.0) which may indicate encryption or packing"
            })
        
        # Check for PE header
        is_pe = header.startswith(b'MZ')
        results["metadata"]["is_pe_format"] = is_pe
        
        if is_pe:
            # Check for digital signature
            has_signature = b'wintrust.dll' in header + footer
            results["metadata"]["appears_signed"] = has_signature
            
            if not has_signature:
                results["indicators"].append({
                    "type": "unsigned_executable",
                    "details": "Executable file appears to be unsigned"
                })
            
            # Extract strings and check for suspicious patterns
            try:
                # Extract ASCII strings (minimum 4 chars)
                strings_ascii = re.findall(b'[ -~]{4,}', header + footer)
                strings_unicode = re.findall(b'(?:[ -~]\x00){4,}', header + footer)
                
                # Convert to text
                strings_list = [s.decode('ascii', errors='ignore') for s in strings_ascii]
                strings_list += [s.decode('utf-16le', errors='ignore') for s in strings_unicode]
                
                # Suspicious API calls and patterns
                suspicious_apis = [
                    "VirtualAlloc", "WriteProcessMemory", "CreateRemoteThread", 
                    "ShellExecute", "WinExec", "CreateProcess", "socket",
                    "InternetOpen", "HttpSendRequest", "UrlDownload",
                    "GetProcAddress", "LoadLibrary", "SetWindowsHook"
                ]
                
                found_apis = []
                for api in suspicious_apis:
                    if any(api in s for s in strings_list):
                        found_apis.append(api)
                
                results["metadata"]["suspicious_apis"] = found_apis
                
                if found_apis:
                    results["indicators"].append({
                        "type": "suspicious_api_imports",
                        "details": f"Executable imports potentially suspicious APIs: {', '.join(found_apis)}"
                    })
                
                # Check for suspicious strings
                suspicious_strings = [
                    "cmd.exe", "powershell", "regsvr32", "rundll32",
                    "http://", "https://", ".onion", ".bit", 
                    "CreateMutex", "Global\\", "\\\\", ".exe",
                    "HKEY_", "SOFTWARE\\", "CurrentVersion\\Run"
                ]
                
                found_strings = []
                for s in suspicious_strings:
                    if any(s.lower() in t.lower() for t in strings_list):
                        found_strings.append(s)
                
                if found_strings:
                    results["indicators"].append({
                        "type": "suspicious_strings",
                        "details": f"Executable contains suspicious strings: {', '.join(found_strings)}"
                    })
                
                # Sample of strings
                results["content_sample"] = "\n".join(strings_list[:50])
                
            except Exception as e:
                results["indicators"].append({
                    "type": "string_extraction_error",
                    "details": f"Error extracting strings: {str(e)}"
                })
    
    except Exception as e:
        results["extraction_error"] = str(e)
    
    return results

def extract_script_content(file_path):
    """Extract and analyze content from script files."""
    extension = os.path.splitext(file_path)[1].lower()
    file_size = os.path.getsize(file_path)
    
    results = {
        "file_type": "script",
        "extension": extension,
        "size_bytes": file_size,
        "metadata": {},
        "indicators": []
    }
    
    try:
        with open(file_path, 'rb') as f:
            content = f.read()
        
        # Try to decode as UTF-8, fallback to Latin-1
        try:
            content_text = content.decode('utf-8')
        except UnicodeDecodeError:
            content_text = content.decode('latin-1')
        
        results["content_sample"] = content_text[:5000]  # First 5000 chars
        
        # Check for obfuscation
        obfuscation_indicators = []
        
        # Base64 patterns
        base64_patterns = re.findall(r'(?:[A-Za-z0-9+/]{4}){8,}(?:[A-Za-z0-9+/]{2}==|[A-Za-z0-9+/]{3}=)?', content_text)
        if base64_patterns:
            obfuscation_indicators.append("base64_encoded_content")
            results["indicators"].append({
                "type": "base64_encoding",
                "details": f"Script contains {len(base64_patterns)} Base64-encoded strings"
            })
        
        # Check for suspicious commands
        suspicious_commands = {
            "powershell -e": "PowerShell Base64 encoded command",
            "powershell -enc": "PowerShell Base64 encoded command",
            "powershell -nop": "PowerShell with execution policy bypass",
            "powershell -w hidden": "PowerShell with hidden window",
            "powershell -exec bypass": "PowerShell execution policy bypass",
            "IEX(": "PowerShell Invoke-Expression (code execution)",
            "Invoke-Expression": "PowerShell code execution",
            "Invoke-WebRequest": "PowerShell web download",
            "wget http": "Download from web",
            "curl http": "Download from web",
            "certutil -decode": "Certutil used for decoding (often malicious)",
            "certutil -urlcache": "Certutil used for downloading",
            "bitsadmin /transfer": "BITSAdmin used for downloading",
            "wscript.shell": "WScript shell object creation",
            "createobject(": "ActiveX object creation",
            "regsvr32 /s /u /i": "Regsvr32 AppLocker bypass",
            "mshta http": "MSHTA remote code execution",
            "rundll32 javascript": "RunDLL32 JavaScript execution",
            "Add-Type": "PowerShell loading C# code",
            "ConvertTo-SecureString": "PowerShell secure string conversion"
        }
        
        found_commands = []
        for cmd, desc in suspicious_commands.items():
            if cmd.lower() in content_text.lower():
                found_commands.append(f"{cmd} ({desc})")
                results["indicators"].append({
                    "type": "suspicious_command",
                    "details": f"Script contains suspicious command: {cmd} - {desc}"
                })
        
        # Check for encoded characters (hex, octal, etc.)
        encoded_char_pattern = re.findall(r'\\x[0-9a-fA-F]{2}|\\u[0-9a-fA-F]{4}|\\[0-7]{3}', content_text)
        if encoded_char_pattern and len(encoded_char_pattern) > 10:
                obfuscation_indicators.append("encoded_chars")
                results["indicators"].append({
                    "type": "character_encoding",
                    "details": f"Script contains {len(encoded_char_pattern)} encoded character sequences (possible obfuscation)"
                })
        
        # Check script type and specific patterns
        if extension in ['.ps1', '.psm1']:
            # PowerShell specific checks
            ps_suspicious = [
                "System.Reflection.Assembly]::Load", "-EncodedCommand", 
                "Get-WmiObject", "Start-Process", "Hidden", "WindowStyle",
                "Net.WebClient", "DownloadString", "DownloadFile",
                "Set-ExecutionPolicy", "Unrestricted", "Bypass",
                "CreateInstance", "DynamicInvoke"
            ]
            
            ps_found = []
            for pattern in ps_suspicious:
                if pattern in content_text:
                    ps_found.append(pattern)
            
            if ps_found:
                results["indicators"].append({
                    "type": "suspicious_powershell",
                    "details": f"PowerShell script uses potentially suspicious methods: {', '.join(ps_found)}"
                })
                
        elif extension == '.bat' or extension == '.cmd':
            # Batch file specific checks
            batch_suspicious = [
                "%COMSPEC%", "call powershell", "call cmd", "start /b",
                "call wscript", "call cscript", "reg add HKCU\\Software\\Microsoft",
                "reg add HKLM\\Software\\Microsoft", "reg add HKCU\\SOFTWARE\\Microsoft\\Windows\\CurrentVersion\\Run"
            ]
            
            batch_found = []
            for pattern in batch_suspicious:
                if pattern.lower() in content_text.lower():
                    batch_found.append(pattern)
            
            if batch_found:
                results["indicators"].append({
                    "type": "suspicious_batch",
                    "details": f"Batch file uses potentially suspicious commands: {', '.join(batch_found)}"
                })
                
        elif extension == '.vbs' or extension == '.js':
            # VBS/JS specific checks
            script_suspicious = [
                "ActiveXObject", "WScript.Shell", "Shell.Application", 
                "Scripting.FileSystemObject", "ADODB.Stream", "Process.Create",
                "Run(", "RegWrite", "GetObject", "eval("
            ]
            
            script_found = []
            for pattern in script_suspicious:
                if pattern in content_text:
                    script_found.append(pattern)
            
            if script_found:
                results["indicators"].append({
                    "type": "suspicious_script",
                    "details": f"Script uses potentially dangerous methods: {', '.join(script_found)}"
                })
            
        # Look for domain names and URLs
        domains = re.findall(r'https?://([a-zA-Z0-9.-]+\.[a-zA-Z]{2,})', content_text)
        results["metadata"]["referenced_domains"] = domains
        
        # Use YARA rules if available
        if yara_rules:
            try:
                matches = yara_rules.match(data=content)
                for match in matches:
                    results["indicators"].append({
                        "type": "yara_match",
                        "details": f"YARA rule matched: {match.rule}"
                    })
            except:
                pass
        
        # Overall obfuscation assessment
        results["metadata"]["obfuscation_indicators"] = obfuscation_indicators
        
    except Exception as e:
        results["extraction_error"] = str(e)
        results["indicators"].append({
            "type": "script_analysis_error",
            "details": f"Error analyzing script: {str(e)}"
        })
    
    return results

def extract_archive_info(file_path):
    """Extract information from archive files (ZIP, RAR, etc.)."""
    extension = os.path.splitext(file_path)[1].lower()
    file_size = os.path.getsize(file_path)
    
    results = {
        "file_type": "archive",
        "extension": extension,
        "size_bytes": file_size,
        "metadata": {},
        "indicators": []
    }
    
    try:
        # For ZIP files
        if extension == '.zip':
            with zipfile.ZipFile(file_path) as zip_ref:
                file_list = zip_ref.namelist()
                
                # Extract metadata
                results["metadata"] = {
                    "file_count": len(file_list),
                    "files": file_list[:100]  # List first 100 files
                }
                
                # Check for suspicious files within archive
                suspicious_extensions = ['.exe', '.dll', '.bat', '.ps1', '.vbs', '.js', '.wsf', '.hta']
                suspicious_files = []
                
                for file in file_list:
                    file_ext = os.path.splitext(file)[1].lower()
                    if file_ext in suspicious_extensions:
                        suspicious_files.append(file)
                
                results["metadata"]["suspicious_files"] = suspicious_files
                
                if suspicious_files:
                    results["indicators"].append({
                        "type": "suspicious_archive_contents",
                        "details": f"Archive contains potentially dangerous files: {', '.join(suspicious_files[:10])}"
                    })
                
                # Check for password protection
                try:
                    test_file = file_list[0] if file_list else None
                    if test_file:
                        zip_ref.read(test_file)
                        results["metadata"]["password_protected"] = False
                except RuntimeError as e:
                    if "password required" in str(e).lower():
                        results["metadata"]["password_protected"] = True
                        results["indicators"].append({
                            "type": "password_protected_archive",
                            "details": "Archive is password protected which may be used to evade detection"
                        })
                
                # Check compression ratio
                total_size = sum(zip_info.file_size for zip_info in zip_ref.infolist())
                compression_ratio = total_size / file_size if file_size > 0 else 0
                results["metadata"]["compression_ratio"] = compression_ratio
                
                if compression_ratio > 100:
                    results["indicators"].append({
                        "type": "zip_bomb",
                        "details": f"Suspicious compression ratio ({compression_ratio:.1f}:1) - possible zip bomb"
                    })
    
    except Exception as e:
        results["extraction_error"] = str(e)
        results["indicators"].append({
            "type": "archive_analysis_error",
            "details": f"Error analyzing archive: {str(e)}"
        })
    
    return results

def extract_apk_content(file_path):
    """Extract and analyze content from Android APK files."""
    import zipfile
    import xml.etree.ElementTree as ET
    import re
    import io
    import os
    import tempfile
    from os.path import basename, splitext, join

    extension = os.path.splitext(file_path)[1].lower()
    file_size = os.path.getsize(file_path)
    
    results = {
        "file_type": "apk",
        "extension": extension,
        "size_bytes": file_size,
        "metadata": {},
        "indicators": [],
        "permissions": [],
        "components": {},
        "content_sample": ""
    }
    
    try:
        # Verify it's a valid APK (which is a ZIP file)
        if not zipfile.is_zipfile(file_path):
            results["indicators"].append({
                "type": "invalid_apk",
                "details": "File has .apk extension but is not a valid ZIP archive"
            })
            return results

        with zipfile.ZipFile(file_path) as apk_zip:
            # Get list of all files in the APK
            file_list = apk_zip.namelist()
            results["metadata"]["file_count"] = len(file_list)
            
            # Look for AndroidManifest.xml
            if 'AndroidManifest.xml' not in file_list:
                results["indicators"].append({
                    "type": "missing_manifest",
                    "details": "APK is missing AndroidManifest.xml"
                })
                return results
            
            # Extract and analyze AndroidManifest.xml
            # Note: The manifest is a binary XML file, we'll extract text content for simple analysis
            manifest_content = apk_zip.read('AndroidManifest.xml')
            
            # Try to extract text strings from manifest to get insights
            strings_from_manifest = re.findall(b'[A-Za-z0-9_\.]{4,}', manifest_content)
            manifest_strings = [s.decode('utf-8', errors='ignore') for s in strings_from_manifest]
            
            # Extract permissions from strings
            permissions = [s for s in manifest_strings if 'android.permission.' in s]
            results["permissions"] = permissions
            
            # Check for dangerous permissions
            dangerous_permissions = [
                "android.permission.SEND_SMS",
                "android.permission.RECEIVE_SMS",
                "android.permission.READ_SMS",
                "android.permission.RECEIVE_WAP_PUSH",
                "android.permission.RECEIVE_MMS",
                "android.permission.READ_PHONE_STATE",
                "android.permission.CALL_PHONE",
                "android.permission.READ_CALL_LOG",
                "android.permission.WRITE_CALL_LOG",
                "android.permission.PROCESS_OUTGOING_CALLS",
                "android.permission.RECORD_AUDIO",
                "android.permission.CAMERA",
                "android.permission.READ_CONTACTS",
                "android.permission.WRITE_CONTACTS",
                "android.permission.READ_CALENDAR",
                "android.permission.WRITE_CALENDAR",
                "android.permission.ACCESS_FINE_LOCATION",
                "android.permission.ACCESS_COARSE_LOCATION",
                "android.permission.ACCESS_BACKGROUND_LOCATION",
                "android.permission.READ_EXTERNAL_STORAGE",
                "android.permission.WRITE_EXTERNAL_STORAGE",
                "android.permission.MANAGE_EXTERNAL_STORAGE",
                "android.permission.REQUEST_INSTALL_PACKAGES",
                "android.permission.SYSTEM_ALERT_WINDOW",
                "android.permission.GET_ACCOUNTS",
                "android.permission.READ_LOGS"
            ]
            
            found_dangerous_permissions = [p for p in permissions if p in dangerous_permissions]
            
            if found_dangerous_permissions:
                results["indicators"].append({
                    "type": "dangerous_permissions",
                    "details": f"APK requests potentially dangerous permissions: {', '.join(found_dangerous_permissions)}"
                })
            
            # Check for obfuscation
            if len([s for s in manifest_strings if len(s) < 3]) > len(manifest_strings) * 0.5:
                results["indicators"].append({
                    "type": "obfuscated_manifest",
                    "details": "Application appears to be obfuscated with short class/method names"
                })
            
            # Look for components
            components = {
                "activities": [s for s in manifest_strings if '.activity.' in s.lower()],
                "services": [s for s in manifest_strings if '.service.' in s.lower()],
                "receivers": [s for s in manifest_strings if '.receiver.' in s.lower()],
                "providers": [s for s in manifest_strings if '.provider.' in s.lower()]
            }
            results["components"] = components
            
            # Check for suspicious component names
            suspicious_component_patterns = [
                "sms", "mms", "message", "receive", "accessibility", "admin", "device",
                "start", "boot", "complete", "screen", "mount", "remote", "shell",
                "command", "exec", "su", "root", "exploit", "crack", "hack"
            ]
            
            suspicious_components = []
            all_components = []
            for comp_type, comps in components.items():
                all_components.extend(comps)
            
            for comp in all_components:
                for pattern in suspicious_component_patterns:
                    if pattern.lower() in comp.lower():
                        suspicious_components.append(comp)
                        break
            
            if suspicious_components:
                results["indicators"].append({
                    "type": "suspicious_components",
                    "details": f"APK contains components with suspicious names: {', '.join(suspicious_components[:5])}"
                })
            
            # Check for admin device privileges
            if 'DeviceAdminReceiver' in '\n'.join(manifest_strings) or 'BIND_DEVICE_ADMIN' in '\n'.join(manifest_strings):
                results["indicators"].append({
                    "type": "device_admin",
                    "details": "Application requests device administrator privileges"
                })
            
            # Extract DEX files count for multi-dex detection
            dex_files = [f for f in file_list if f.endswith('.dex')]
            results["metadata"]["dex_files"] = dex_files
            
            if len(dex_files) > 1:
                results["metadata"]["multi_dex"] = True
            
            # Look for encrypted/packed code (detect libs that might be used for packing)
            native_libs = [f for f in file_list if f.startswith('lib/') and f.endswith('.so')]
            results["metadata"]["native_libraries"] = [basename(lib) for lib in native_libs]
            
            # Check for known packer signatures
            packer_libs = [
                "libsecexe.so", "libprotectClass.so", "libDexHelper.so", "libshell.so",
                "libjiagu.so", "libBugly.so", "libmobisec.so", "libpreverify.so"
            ]
            
            found_packer_libs = []
            for lib in native_libs:
                lib_name = basename(lib)
                if any(packer in lib_name for packer in packer_libs):
                    found_packer_libs.append(lib_name)
            
            if found_packer_libs:
                results["indicators"].append({
                    "type": "app_packer_detected",
                    "details": f"APK appears to be using code packing/protection: {', '.join(found_packer_libs)}"
                })
            
            # Check for presence of embedded APKs or ZIPs (could be used for dynamic loading)
            embedded_apks = [f for f in file_list if f.endswith('.apk') or (f.endswith('.zip') and not f.startswith('META-INF/'))]
            
            if embedded_apks:
                results["indicators"].append({
                    "type": "embedded_apk",
                    "details": f"APK contains embedded APK or ZIP files which may be used for dynamic loading: {', '.join(embedded_apks)}"
                })
            
            # Extract and analyze DEX for suspicious patterns
            if dex_files:
                try:
                    dex_content = apk_zip.read(dex_files[0])  # Analyze the first DEX file
                    
                    # Look for suspicious strings
                    dex_strings = re.findall(b'[A-Za-z0-9_\.:/]{4,}', dex_content)
                    dex_text = [s.decode('utf-8', errors='ignore') for s in dex_strings]
                    
                    # Extract a sample of strings for content preview
                    results["content_sample"] = "\n".join(dex_text[:100])
                    
                    # Check for known suspicious API calls
                    suspicious_api_patterns = [
                        "Landroid/telephony/SmsManager;->send",
                        "Ljava/lang/Runtime;->exec",
                        "Landroid/app/admin/DevicePolicy",
                        "Ljava/lang/reflect/Method;->invoke",
                        "Ldalvik/system/DexClassLoader",
                        "Ldalvik/system/PathClassLoader",
                        "Ljava/net/URLClassLoader",
                        "Lorg/apache/http/impl/client/DefaultHttpClient",
                        "Ljava/net/HttpURLConnection",
                        "Landroid/content/pm/PackageManager;->setComponentEnabledSetting",
                        "Landroid/app/Service;->startForeground",
                        "Landroid/app/NotificationManager",
                        "Landroid/content/Context;->startService",
                        "Landroid/location/LocationManager",
                        "Landroid/hardware/Camera",
                        "Landroid/media/AudioRecord"
                    ]
                    
                    found_suspicious_apis = []
                    dex_text_joined = "\n".join(dex_text)
                    
                    for api in suspicious_api_patterns:
                        if api in dex_text_joined:
                            found_suspicious_apis.append(api)
                    
                    if found_suspicious_apis:
                        results["indicators"].append({
                            "type": "suspicious_api_usage",
                            "details": f"DEX contains calls to potentially dangerous APIs: {', '.join(found_suspicious_apis[:5])}"
                        })
                    
                    # Check for reflection & dynamic code
                    reflection_patterns = [
                        "Ljava/lang/reflect",
                        "Ljava/lang/Class;->forName",
                        "Ljava/lang/Class;->getMethod",
                        "Ljava/lang/reflect/Method",
                        "Ljavassist",
                        "Ldalvik/system/DexFile",
                        "Ldalvik/system/InMemoryDexClassLoader"
                    ]
                    
                    reflection_apis = [api for api in reflection_patterns if api in dex_text_joined]
                    
                    if reflection_apis:
                        results["indicators"].append({
                            "type": "reflection_usage",
                            "details": f"DEX contains code using reflection which can be used to hide functionality: {', '.join(reflection_apis)}"
                        })
                    
                    # Check for encrypted strings or native method usage
                    native_method_calls = len(re.findall(b'native', dex_content))
                    if native_method_calls > 5:
                        results["indicators"].append({
                            "type": "heavy_native_code",
                            "details": f"DEX contains multiple native method declarations ({native_method_calls}) which can be used to hide functionality"
                        })
                    
                    # Check for common C2 domains
                    domains = re.findall(r'https?://([a-zA-Z0-9.-]+\.[a-zA-Z]{2,})', dex_text_joined)
                    for domain in domains:
                        # Check for dynamic DNS services often used by malware
                        dynamic_dns = [
                            "dyndns.org", "noip.com", "ddns.net", "serveo.net", "ngrok.io", 
                            "hopto.org", "zapto.org", "3utilities.com", "bounceme.net", "freedynamicdns.net",
                            "freedynamicdns.org", "gotdns.ch", "hopto.org", "myddns.me", "nsupdate.info", 
                            "servebeer.com", "serveblog.net", "serveminecraft.net", "sytes.net", "zapto.org"
                        ]
                        
                        if any(d in domain.lower() for d in dynamic_dns):
                            results["indicators"].append({
                                "type": "dynamic_dns_usage",
                                "details": f"App communicates with dynamic DNS service: {domain}"
                            })
                    
                    # Look for command and control servers or suspicious URLs
                    c2_patterns = [
                        "/command", "/control", "/bot", "/admin", "/gate", "/gate.php", 
                        "/panel", "/server", "/client", "/config", "/cfg", "/settings"
                    ]
                    
                    for url in re.findall(r'https?://[a-zA-Z0-9.-/]+', dex_text_joined):
                        if any(pattern in url.lower() for pattern in c2_patterns):
                            results["indicators"].append({
                                "type": "potential_c2_url",
                                "details": f"Suspicious URL possibly used for command & control: {url}"
                            })
                        
                except Exception as e:
                    results["indicators"].append({
                        "type": "dex_analysis_error",
                        "details": f"Error analyzing DEX file: {str(e)}"
                    })
            
            # Look for encrypted assets
            asset_files = [f for f in file_list if f.startswith('assets/')]
            if asset_files:
                # Try to detect encrypted assets by checking entropy
                encrypted_assets = []
                for asset in asset_files[:5]:  # Check only first 5 assets
                    try:
                        asset_content = apk_zip.read(asset)
                        entropy = get_file_entropy(asset_content)
                        if entropy > 7.8:  # Very high entropy suggests encryption
                            encrypted_assets.append(f"{asset} (entropy: {entropy:.2f})")
                    except:
                        pass
                
                if encrypted_assets:
                    results["indicators"].append({
                        "type": "encrypted_assets",
                        "details": f"APK contains potentially encrypted assets: {', '.join(encrypted_assets)}"
                    })
            
            # Check for root detection or anti-emulation techniques
            root_detection_patterns = [
                "su", "/system/bin/su", "/system/xbin/su", "/sbin/su", 
                "Superuser.apk", "supersu", "chainfire", "busybox",
                "rootcloak", "rootbeer"
            ]
            
            anti_emulation_patterns = [
                "android.os.Build", "FINGERPRINT", "BRAND", "DEVICE", "MANUFACTURER", "MODEL",
                "HARDWARE", "isEmulator", "qemu", "generic", "goldfish", "sdk", "vbox",
                "virtual"
            ]
            
            root_checks = []
            emulator_checks = []
            
            for pattern in root_detection_patterns:
                if pattern.lower() in dex_text_joined.lower():
                    root_checks.append(pattern)
            
            for pattern in anti_emulation_patterns:
                # Count occurrences to avoid false positives
                if pattern.lower() in dex_text_joined.lower() and dex_text_joined.lower().count(pattern.lower()) > 2:
                    emulator_checks.append(pattern)
            
            if root_checks:
                results["metadata"]["root_detection"] = True
                
            if emulator_checks and len(emulator_checks) > 3:  # Require multiple matches to reduce false positives
                results["metadata"]["anti_emulation"] = True
                results["indicators"].append({
                    "type": "anti_analysis_techniques",
                    "details": f"APK contains code that appears to detect root/emulators, potentially to evade analysis"
                })
            
            # Check for ad frameworks and SDK identifiers
            ad_frameworks = [
                "com.google.android.gms.ads", "com.mopub", "com.facebook.ads",
                "com.unity3d.ads", "com.applovin", "com.chartboost", "com.vungle",
                "com.tapjoy", "com.ironsource", "com.flurry", "com.adcolony",
                "com.startapp", "com.inmobi", "com.appbrain"
            ]
            
            found_ad_sdks = []
            for sdk in ad_frameworks:
                if sdk in dex_text_joined:
                    found_ad_sdks.append(sdk)
            
            if found_ad_sdks:
                results["metadata"]["ad_frameworks"] = found_ad_sdks
            
            # Check for potentially unwanted third-party SDKs
            tracking_sdks = [
                "com.singular.sdk", "com.amplitude", "com.appsflyer", "com.kochava",
                "io.branch", "com.adjust", "com.localytics", "com.google.firebase.analytics",
                "com.mixpanel", "com.segment", "io.fabric"
            ]
            
            found_tracking_sdks = []
            for sdk in tracking_sdks:
                if sdk in dex_text_joined:
                    found_tracking_sdks.append(sdk)
            
            if found_tracking_sdks:
                results["metadata"]["tracking_sdks"] = found_tracking_sdks
    
    except Exception as e:
        results["extraction_error"] = str(e)
        results["indicators"].append({
            "type": "apk_analysis_error",
            "details": f"Error analyzing APK: {str(e)}"
        })
    
    return results

def extract_file_content(file_path):
    """Extract content from different file types for analysis."""
    file_name = os.path.basename(file_path)
    extension = os.path.splitext(file_path)[1].lower()
    
    # Get file type using magic
    try:
        import magic
        mime_type = magic.Magic(mime=True).from_file(file_path)
        file_type_desc = magic.Magic().from_file(file_path)
    except ImportError:
        mime_type = mimetypes.guess_type(file_path)[0] or 'unknown/unknown'
        file_type_desc = f"Based on extension: {extension}"
    
    # Calculate file hashes
    hashes = get_file_hashes(file_path)
    
    # For APK files
    if extension == '.apk' or mime_type == 'application/vnd.android.package-archive':
        result = extract_apk_content(file_path)
        result["file_name"] = file_name
        result["mime_type"] = mime_type
        result["file_desc"] = file_type_desc
        result["hashes"] = hashes
        return result
    
    # For PDF files
    elif extension == '.pdf' or mime_type == 'application/pdf':
        result = extract_pdf_content(file_path)
        result["file_name"] = file_name
        result["mime_type"] = mime_type
        result["file_desc"] = file_type_desc
        result["hashes"] = hashes
        return result
    
    # For Office documents, use the specialized function
    elif extension in ['.docx', '.xlsx', '.doc', '.xls', '.rtf', '.ppt', '.pptx']:
        result = extract_office_document_content(file_path)
        result["file_name"] = file_name
        result["mime_type"] = mime_type
        result["file_desc"] = file_type_desc
        result["hashes"] = hashes
        return result
    
    # For binary files, get metadata and hex dump
    elif extension in ['.exe', '.dll', '.sys']:
        result = extract_executable_info(file_path)
        result["file_name"] = file_name
        result["mime_type"] = mime_type
        result["file_desc"] = file_type_desc
        result["hashes"] = hashes
        return result
    
    # For script files
    elif extension in ['.bat', '.ps1', '.vbs', '.js', '.wsf', '.hta', '.cmd']:
        result = extract_script_content(file_path)
        result["file_name"] = file_name
        result["mime_type"] = mime_type
        result["file_desc"] = file_type_desc
        result["hashes"] = hashes
        return result
    
    # For archive files
    elif extension in ['.zip', '.rar', '.7z', '.tar', '.gz']:
        result = extract_archive_info(file_path)
        result["file_name"] = file_name
        result["mime_type"] = mime_type
        result["file_desc"] = file_type_desc
        result["hashes"] = hashes
        return result
    
    # Default for other files
    else:
        try:
            with open(file_path, 'rb') as f:
                content = f.read(10000)
            
            # Try to detect embedded files
            embedded_files = detect_embedded_files(file_path)
            
            # Try to decode as text
            try:
                content_text = content.decode('utf-8')
                file_type = "text"
            except UnicodeDecodeError:
                content_text = "Binary content"
                file_type = "binary"
            
            result = {
                "file_name": file_name,
                "file_type": file_type,
                "extension": extension,
                "mime_type": mime_type,
                "file_desc": file_type_desc,
                "content": content_text if file_type == "text" else "Binary content",
                "size_bytes": os.path.getsize(file_path),
                "hashes": hashes,
                "metadata": {},
                "indicators": []
            }
            
            # Add embedded files as an indicator if found
            if embedded_files:
                result["metadata"]["embedded_files"] = embedded_files
                result["indicators"].append({
                    "type": "embedded_files",
                    "details": f"File contains signatures of embedded content: {', '.join([e['type'] for e in embedded_files])}"
                })
            
            # Use YARA rules if available
            if yara_rules and file_type == "binary":
                try:
                    matches = yara_rules.match(data=content)
                    for match in matches:
                        result["indicators"].append({
                            "type": "yara_match",
                            "details": f"YARA rule matched: {match.rule}"
                        })
                except:
                    pass
            
            return result
            
        except Exception as e:
            return {
                "file_name": file_name,
                "file_type": "unknown",
                "extension": extension,
                "size_bytes": os.path.getsize(file_path),
                "hashes": hashes,
                "error": str(e)
            }

analysis_template = ChatPromptTemplate.from_messages([
    ("system", """
    You are a cybersecurity expert specializing in malware analysis. Analyze the provided file content or metadata 
    to determine if it contains malicious patterns. Consider all the data provided, including:
    
    1. File metadata (type, size, extension)
    2. Extracted content samples
    3. Detected indicators of potentially malicious patterns
    4. Hash values
    5. Embedded objects or files
    6. Suspicious strings or commands
    
    GUIDELINES FOR ANALYSIS:
    
    - PDF Files: Look for JavaScript, embedded files, suspicious URI actions, exploits
    - Office Documents: Macros, embedded objects, suspicious links, DDE commands, OLE objects
    - Scripts (.bat, .ps1, .vbs): Obfuscation, suspicious commands, encoded payloads
    - Executables: High entropy (packing), suspicious imports, lack of signatures
    - Archives: Suspicious compression ratios, password protection, malicious content
    - APK Files: Dangerous permissions, suspicious components, anti-analysis techniques, embedded payloads,
      suspicious API calls, device admin requests, reflection usage, native code, C2 communication
    
    For APK files specifically, consider these risk factors:
    - Numerous dangerous permissions (SMS, CALL_PHONE, ADMIN)
    - Use of reflection or dynamic code loading
    - Presence of code obfuscation or packing
    - Anti-emulation or root detection techniques (could be evasion)
    - Embedded APKs or encrypted assets
    - Suspicious URLs or communication patterns
    - Device administrator privileges requests
    
    Evaluate each indicator carefully. Consider both the presence of suspicious elements AND context.
    
    RESPONSE FORMAT INSTRUCTIONS:
    Return ONLY a JSON object with this exact structure and nothing else:
    {{
        "file_name": "name of the file",
        "file_type": "type of file",
        "is_malicious": true/false,
        "confidence": 0-100,
        "risk_level": "low/medium/high/critical",
        "indicators": [
            {{"type": "description of indicator", "details": "details about why this is suspicious", "severity": "low/medium/high"}}
        ],
        "mitigations": [
            "Suggestion for handling this file safely"
        ],
        "analysis_summary": "A paragraph summarizing your analysis and conclusion"
    }}
    
    DO NOT include any text, explanations, or markdown before or after the JSON object.
    DO NOT use ```json or ``` markup around the JSON.
    The response should be ONLY the JSON object itself.
    """),
    ("human", "{file_info}")
])

analysis_chain = analysis_template | model | StrOutputParser()

def extract_json_from_text(text):
    """Extract a JSON object from text that might contain additional content."""
    json_match = re.search(r'({[\s\S]*})', text)
    
    if json_match:
        json_str = json_match.group(1)
        try:
            return json.loads(json_str)
        except json.JSONDecodeError:
            pass
    
    lines = []
    in_json = False
    for line in text.split('\n'):
        line = line.strip()
        if line == '{':
            in_json = True
            lines.append(line)
        elif line == '}':
            lines.append(line)
            in_json = False
        elif in_json or (line.startswith('"') and ':' in line):
            lines.append(line)
    
    if lines:
        try:
            return json.loads('\n'.join(lines))
        except json.JSONDecodeError:
            pass
    
    return None

def estimate_tokens(text):
    # Rough estimate: 1 token ≈ 4 characters (varies by language/model)
    return len(text) / 4

def truncate_for_token_limit(file_content, max_tokens=4000):
    """Truncate file content to stay under token limit while preserving important information."""
    # Estimate function
    def estimate_tokens(text):
        # More accurate estimate: ~4 chars per token for English text
        return len(json.dumps(text)) / 4
    
    # Keep essential metadata
    essential = {
        "file_name": file_content.get("file_name", ""),
        "file_type": file_content.get("file_type", ""),
        "extension": file_content.get("extension", ""),
        "mime_type": file_content.get("mime_type", ""),
        "size_bytes": file_content.get("size_bytes", 0),
        "hashes": file_content.get("hashes", {})
    }

    # Start with most critical information
    truncated = {
        **essential,
        "indicators": [],
        "metadata": {},
        "content_sample": ""
    }
    
    # Reserve tokens for structure overhead
    structure_overhead = 100
    available_tokens = max_tokens - structure_overhead
    
    # Prioritize indicators based on severity/importance
    indicators = sorted(
        file_content.get("indicators", []),
        key=lambda x: 0 if x.get("type", "").startswith(("malicious", "exploit", "suspicious")) else 1
    )
    
    # Calculate initial token usage
    tokens_used = estimate_tokens(truncated)
    
    # Add indicators - most important signals of malicious content
    for indicator in indicators:
        indicator_copy = {
            "type": indicator.get("type", ""),
            "details": indicator.get("details", "")[:200]  # Limit details length
        }
        
        # Check if adding this would exceed our limit
        with_indicator = {**truncated, "indicators": truncated["indicators"] + [indicator_copy]}
        new_tokens = estimate_tokens(with_indicator)
        
        if new_tokens <= available_tokens:
            truncated["indicators"].append(indicator_copy)
            tokens_used = new_tokens
        else:
            # If we can't add more indicators, move on
            break
    
    # Add critical metadata fields (prioritize security-relevant fields)
    priority_metadata = [
        "has_javascript", "has_macros", "has_vba_macros", "has_embedded_files", 
        "embedded_objects", "suspicious_apis", "obfuscation_indicators",
        "has_form_fields", "entropy", "password_protected"
    ]
    
    # First add priority metadata
    for key in priority_metadata:
        if key in file_content.get("metadata", {}):
            temp = {**truncated, "metadata": {**truncated["metadata"], key: file_content["metadata"][key]}}
            new_tokens = estimate_tokens(temp)
            
            if new_tokens <= available_tokens:
                truncated["metadata"][key] = file_content["metadata"][key]
                tokens_used = new_tokens
            else:
                break
    
    # Then add other metadata if space permits
    for key, value in file_content.get("metadata", {}).items():
        if key in truncated["metadata"]:
            continue  # Skip already added priority metadata
            
        # Handle lists specially - truncate long lists
        if isinstance(value, list) and len(value) > 5:
            shortened_value = value[:5]  # Only keep first 5 items
            temp = {**truncated, "metadata": {**truncated["metadata"], key: shortened_value}}
        else:
            temp = {**truncated, "metadata": {**truncated["metadata"], key: value}}
            
        new_tokens = estimate_tokens(temp)
        if new_tokens <= available_tokens:
            if isinstance(value, list) and len(value) > 5:
                truncated["metadata"][key] = value[:5]
            else:
                truncated["metadata"][key] = value
            tokens_used = new_tokens
        else:
            break
    
    # Add content sample with remaining token budget
    remaining_tokens = available_tokens - tokens_used
    content = file_content.get("content_sample", "")
    
    # Estimate how much content we can include
    chars_per_token = 4
    content_char_limit = int(remaining_tokens * chars_per_token)
    
    # Ensure we don't go over budget
    if content and content_char_limit > 20:  # Only add if we can include a meaningful amount
        truncated["content_sample"] = content[:content_char_limit]
    
    # Final check to ensure we're under the limit
    final_tokens = estimate_tokens(truncated)
    if final_tokens > max_tokens:
        # Emergency truncation if still over limit
        if "content_sample" in truncated:
            truncated["content_sample"] = truncated["content_sample"][:int(len(truncated["content_sample"])*0.7)]
        
        # If still over, truncate indicators
        if estimate_tokens(truncated) > max_tokens and truncated["indicators"]:
            truncated["indicators"] = truncated["indicators"][:max(1, len(truncated["indicators"])//2)]
    
    return truncated

def analyze_file_for_malware(file_path):
    """Process a file and determine if it might contain malware."""
    file_content = extract_file_content(file_path)
    logger.info(f"Analyzing file: {file_content['file_name']}")
    
    try:
        # Apply token limiting
        truncated_content = truncate_for_token_limit(file_content)
        
        analysis_result = analysis_chain.invoke({"file_info": json.dumps(truncated_content, indent=2)})
        
        try:
            result_json = json.loads(analysis_result)
            return result_json
        except json.JSONDecodeError:
            extracted_json = extract_json_from_text(analysis_result)
            
            if extracted_json:
                return extracted_json
            else:
                return {
                    "file_name": file_content["file_name"],
                    "is_malicious": "unknown",
                    "error": "Could not parse analysis result"
                }
            
    except Exception as e:
        logger.error(f"Error during analysis: {str(e)}")
        return {
            "file_name": file_content["file_name"],
            "is_malicious": "unknown",
            "error": str(e)
        }

def analyze_directory(directory_path, output_file=None):
    """Analyze all files in a directory and output results."""
    results = []
    
    for root, _, files in os.walk(directory_path):
        for file_name in files:
            file_path = os.path.join(root, file_name)
            logger.info(f"Processing {file_path}")
            
            try:
                result = analyze_file_for_malware(file_path)
                results.append(result)
                
                # Print result
                print(f"File: {result['file_name']}")
                print(f"Analysis: {'MALICIOUS' if result.get('is_malicious') else 'Clean'}")
                print(f"Confidence: {result.get('confidence', 'N/A')}")
                print(f"Risk Level: {result.get('risk_level', 'N/A')}")
                print("-" * 50)
                
            except Exception as e:
                logger.error(f"Failed to analyze {file_path}: {str(e)}")
    
    # Write results to file if requested
    if output_file:
        with open(output_file, 'w') as f:
            json.dump(results, f, indent=2)
    
    return results